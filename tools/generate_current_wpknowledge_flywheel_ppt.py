#!/usr/bin/env python3
"""Generate the editable project-report deck for wpKnowledge PR #20.

Dependency: python-pptx >= 1.0
Run from the repository root:
  python tools/generate_current_wpknowledge_flywheel_ppt.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "knowledge/2.wiki/设计/当前wpKnowledge知识飞轮方案.pptx"

W = 13.333
H = 7.5
FONT = "Microsoft YaHei"
MONO = "Aptos Mono"

NAVY = "10253F"
NAVY_2 = "173A5E"
INK = "162231"
MUTED = "607084"
BG = "F4F7FA"
WHITE = "FFFFFF"
CYAN = "16B7C9"
CYAN_LIGHT = "DDF6F8"
BLUE = "3976E8"
BLUE_LIGHT = "E8F0FF"
GREEN = "229B6A"
GREEN_LIGHT = "DFF3E9"
ORANGE = "E8892E"
ORANGE_LIGHT = "FFF0DC"
RED = "D95852"
RED_LIGHT = "FBE7E5"
PURPLE = "7656D6"
PURPLE_LIGHT = "EDE8FB"
LINE = "D7E0E8"


def rgb(hex_value: str) -> RGBColor:
    return RGBColor.from_string(hex_value)


prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
prs.core_properties.title = "当前wpKnowledge知识飞轮方案"
prs.core_properties.subject = "wpKnowledge PR #20 代码评审与知识飞轮方案汇报"
prs.core_properties.author = "wpKnowledge project team"


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def textbox(slide, x, y, w, h, text, size=16, color=INK, bold=False,
            align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE, font=FONT,
            margin=0.04, fit=True, italic=False):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    tf.word_wrap = True
    if fit:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = rgb(color)
    return shape


def box(slide, x, y, w, h, text="", fill=WHITE, line=LINE, radius=True,
        size=15, color=INK, bold=False, align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE, line_width=1.0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(line_width)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.10)
    tf.margin_right = Inches(0.10)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    tf.vertical_anchor = valign
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)
    return shape


def rich_box(slide, x, y, w, h, title, lines, fill=WHITE, accent=CYAN,
             title_color=INK, body_color=MUTED):
    box(slide, x, y, w, h, fill=fill, line=LINE, text="")
    box(slide, x, y, 0.08, h, fill=accent, line=accent, radius=False, text="")
    textbox(slide, x + 0.22, y + 0.12, w - 0.34, 0.38, title, 17, title_color, True)
    body = "\n".join(f"• {line}" for line in lines)
    textbox(slide, x + 0.22, y + 0.55, w - 0.34, h - 0.67, body, 12.5, body_color,
            valign=MSO_ANCHOR.TOP)


def pill(slide, x, y, w, text, fill, color=WHITE, size=10.5):
    return box(slide, x, y, w, 0.29, text, fill=fill, line=fill, size=size,
               color=color, bold=True)


def chevron(slide, x, y, w=0.30, h=0.42, color=CYAN):
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.color.rgb = rgb(color)
    return shape


def line(slide, x1, y1, x2, y2, color=LINE, width=1.5, dash=None):
    shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width)
    if dash is not None:
        shape.line.dash_style = dash
    return shape


def slide_frame(title, status, source, dark=False, kicker=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY if dark else BG)
    title_color = WHITE if dark else NAVY
    muted_color = "B9C9D8" if dark else MUTED
    if kicker:
        textbox(slide, 0.55, 0.30, 4.5, 0.24, kicker.upper(), 9.5, CYAN, True)
    textbox(slide, 0.55, 0.52, 10.8, 0.55, title, 25, title_color, True)
    status_color = {"IMPLEMENTED": GREEN, "DEMO": ORANGE, "PLANNED": PURPLE,
                    "IMPLEMENTED + DEMO": BLUE}.get(status, MUTED)
    pill(slide, 11.42, 0.45, 1.35, status, status_color, size=9.0)
    line(slide, 0.55, 1.18, 12.78, 1.18, "31516E" if dark else LINE, 1.0)
    textbox(slide, 0.55, 7.13, 11.75, 0.20, source, 8.0, muted_color,
            valign=MSO_ANCHOR.BOTTOM, fit=False)
    textbox(slide, 12.35, 7.10, 0.43, 0.22, str(len(prs.slides)), 8.5, muted_color,
            align=PP_ALIGN.RIGHT, fit=False)
    return slide


def node(slide, x, y, w, h, title, subtitle="", fill=WHITE, accent=CYAN,
         title_color=INK, body_color=MUTED):
    box(slide, x, y, w, h, fill=fill, line=accent, text="", line_width=1.3)
    textbox(slide, x + 0.12, y + 0.08, w - 0.24, 0.30, title, 14, title_color, True,
            align=PP_ALIGN.CENTER)
    if subtitle:
        textbox(slide, x + 0.12, y + 0.40, w - 0.24, h - 0.47, subtitle, 10.5,
                body_color, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.TOP)


def status_legend(slide, y=6.63, dark=False):
    base = "B9C9D8" if dark else MUTED
    textbox(slide, 8.12, y, 0.75, 0.25, "状态口径", 9, base, True)
    pill(slide, 8.90, y - 0.01, 1.05, "IMPLEMENTED", GREEN, size=8.0)
    pill(slide, 10.06, y - 0.01, 0.68, "DEMO", ORANGE, size=8.0)
    pill(slide, 10.85, y - 0.01, 0.90, "PLANNED", PURPLE, size=8.0)
    textbox(slide, 11.84, y, 0.90, 0.25, "以代码为准", 8.5, base,
            align=PP_ALIGN.RIGHT)


# 1 — title
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NAVY)
box(slide, 0, 0, 0.18, H, fill=CYAN, line=CYAN, radius=False)
textbox(slide, 0.72, 0.78, 4.0, 0.28, "WPKNOWLEDGE · PR #20", 11, CYAN, True)
textbox(slide, 0.72, 1.25, 11.8, 1.15, "当前wpKnowledge知识飞轮方案", 32, WHITE, True,
        valign=MSO_ANCHOR.TOP)
textbox(slide, 0.76, 2.72, 8.8, 0.42,
        "代码评审结论 · 多 Agent demo 进度 · 语义分块方案 · 下一步", 16, "C5D3E0")
box(slide, 0.75, 3.55, 11.85, 1.55, fill=NAVY_2, line="31516E", text="")
textbox(slide, 1.05, 3.84, 3.0, 0.28, "本次基线", 10, CYAN, True)
textbox(slide, 1.05, 4.16, 3.0, 0.48, "PR #20\n776dac5 → b6973a8", 15, WHITE, True,
        valign=MSO_ANCHOR.TOP)
textbox(slide, 4.62, 3.84, 3.0, 0.28, "验证结果", 10, CYAN, True)
textbox(slide, 4.62, 4.16, 3.0, 0.48, "typecheck / specs\n50 / 50 tests", 15, WHITE, True,
        valign=MSO_ANCHOR.TOP)
textbox(slide, 8.18, 3.84, 3.0, 0.28, "汇报日期", 10, CYAN, True)
textbox(slide, 8.18, 4.16, 3.0, 0.48, "2026-09-02\nAsia / Beijing", 15, WHITE, True,
        valign=MSO_ANCHOR.TOP)
textbox(slide, 0.75, 6.68, 11.2, 0.22,
        "范围：wpKnowledge 当前分支 codex/embedded-domain-knowledge；domain-knowledge 仅作只读对照。",
        9, "9FB2C5", fit=False)
pill(slide, 10.42, 6.62, 1.05, "IMPLEMENTED", GREEN, size=8.0)
pill(slide, 11.58, 6.62, 0.68, "DEMO", ORANGE, size=8.0)
pill(slide, 12.34, 6.62, 0.80, "PLANNED", PURPLE, size=8.0)


# 2 — review conclusion
slide = slide_frame(
    "评审结论：边界方向正确，两处事实源问题已修复",
    "IMPLEMENTED",
    "代码依据：packages/application/src/automated-project-workflow.ts；packages/domain/src/index.ts；commit b6973a8",
    kicker="CODE REVIEW",
)
rich_box(slide, 0.60, 1.48, 3.82, 2.03, "边界合理", [
    "LangGraph 只负责节点调度、重试、取消与 checkpoint",
    "wpKnowledge 保留 Run、版本、证据、Gate 与 Publication",
    "Infrastructure → Ports → Application / Domain 的依赖方向未反转",
], accent=GREEN)
rich_box(slide, 4.75, 1.48, 3.82, 2.03, "修复 01 · Gate 顺序", [
    "原实现先产生 GateDecision，再运行 Review",
    "现由 Review 完成后绑定 Oracle / Check / Review 产物",
    "Check 或 Review blocking 会得到 ITERATE / STOPPED",
], accent=RED)
rich_box(slide, 8.90, 1.48, 3.82, 2.03, "修复 02 · 单一事实源", [
    "GraphState 与 FlywheelRun 均从 iteration 0 起步",
    "workflow_router 不再把业务 ITERATE 私自改为 STOPPED",
    "只有 Domain Gate 决定 PASS / ITERATE / STOPPED",
], accent=RED)
box(slide, 0.60, 3.86, 12.12, 1.77, fill=NAVY, line=NAVY, text="")
textbox(slide, 0.90, 4.10, 2.2, 0.31, "最终权威链", 12, CYAN, True)
steps = [
    ("CANDIDATE", BLUE), ("EvaluationReport", PURPLE), ("GateDecision", ORANGE),
    ("wp Publication", GREEN), ("VERIFIED", GREEN),
]
x = 0.90
for index, (label, color) in enumerate(steps):
    box(slide, x, 4.53, 1.75 if index else 1.45, 0.55, label, fill=NAVY_2,
        line=color, color=WHITE, size=12, bold=True)
    x += 1.93 if index else 1.65
    if index < len(steps) - 1:
        chevron(slide, x - 0.14, 4.60, 0.28, 0.38, CYAN)
textbox(slide, 10.02, 4.23, 2.35, 0.73,
        "只有最后一步\n可以写入 VERIFIED", 15, WHITE, True, align=PP_ALIGN.CENTER)
status_legend(slide)


# 3 — boundary
slide = slide_frame(
    "职责边界：执行状态与业务事实并行，但不互相冒充",
    "IMPLEMENTED",
    "代码依据：infrastructure/domain-knowledge/src/{graph,state,runtime}.ts；apps/runner/src/composition.ts:72-129",
    dark=True,
    kicker="ARCHITECTURE",
)
textbox(slide, 0.63, 1.42, 2.5, 0.25, "LANGGRAPH · 执行控制", 11, CYAN, True)
box(slide, 0.60, 1.79, 5.78, 3.98, fill=NAVY_2, line="31516E", text="")
for i, (title, sub, color) in enumerate([
    ("GraphState", "currentNode · route · iteration", BLUE),
    ("Checkpoint", "thread_id · resume · attempts", PURPLE),
    ("Projection", "RUNNING / COMPLETED / FAILED", CYAN),
    ("AbortSignal", "协作式取消", ORANGE),
]):
    node(slide, 0.92 + (i % 2) * 2.60, 2.18 + (i // 2) * 1.38, 2.25, 0.92,
         title, sub, fill=NAVY, accent=color, title_color=WHITE, body_color="B9C9D8")
textbox(slide, 0.92, 5.18, 5.05, 0.34, "不持有 KnowledgeVersion / GateDecision / Publication", 12,
        "C5D3E0", True, align=PP_ALIGN.CENTER)

textbox(slide, 6.94, 1.42, 3.3, 0.25, "WPKNOWLEDGE REGISTRY · 业务事实", 11, CYAN, True)
box(slide, 6.90, 1.79, 5.82, 3.98, fill="F7FAFC", line="31516E", text="")
for i, (title, sub, color) in enumerate([
    ("FlywheelRun", "状态 + 业务 iteration", BLUE),
    ("CAS / Evidence", "sha256 不可变引用", PURPLE),
    ("GateDecision", "策略 + reasonCodes", ORANGE),
    ("Publication", "原子写入 VERIFIED", GREEN),
]):
    node(slide, 7.22 + (i % 2) * 2.60, 2.18 + (i // 2) * 1.38, 2.25, 0.92,
         title, sub, fill=WHITE, accent=color)
textbox(slide, 7.24, 5.17, 5.05, 0.35, "业务事件与 SQLite 投影由同一 Application Service 写入", 12,
        INK, True, align=PP_ALIGN.CENTER)
chevron(slide, 6.48, 3.40, 0.30, 0.48, CYAN)
textbox(slide, 5.90, 3.90, 1.45, 0.45, "Ports / contracts", 9.5, "B9C9D8", True,
        align=PP_ALIGN.CENTER)
status_legend(slide, dark=True)


# 4 — progress matrix
slide = slide_frame(
    "当前进度：编排骨架已落地，真实 Agent 与语义分块仍未接入",
    "IMPLEMENTED + DEMO",
    "代码依据：contracts/src/index.ts:200-216；composition.ts:97-114；tests/acceptance/automated-langgraph-flow.test.ts",
    kicker="PROGRESS",
)
headers = [("能力", 0.65, 4.35), ("IMPLEMENTED", 5.02, 1.68), ("DEMO", 6.74, 1.32),
           ("PLANNED", 8.10, 1.46), ("判断", 9.60, 3.05)]
for label, x, w in headers:
    box(slide, x, 1.46, w, 0.43, label, fill=NAVY, line=NAVY, color=WHITE, size=10.5, bold=True)
rows = [
    ("7 Agent 节点 + LangGraph 路由", True, False, False, "节点、投影、fan-out/join 可运行"),
    ("SQLite graph checkpoint + Registry/CAS", True, False, False, "两类状态分层持久化"),
    ("ohMyWorkPanel 两轮闭环", False, True, False, "静态资产驱动，不是 live Agent"),
    ("TrustedProjectEvaluator", True, True, False, "真实命令执行；不是敌对沙箱"),
    ("promptAddon-only 配置", True, False, False, "字段边界安全；执行端尚不消费"),
    ("真实 Agent / CodeAgent provider", False, False, True, "自动路径尚未接入"),
    ("拓扑 + 语义切块与增量上下文", False, False, True, "当前 worker 只是机械 fixture"),
    ("进程级取消 / RUNNING 恢复", False, False, True, "当前只证明协作式取消、部分恢复"),
]
for idx, row in enumerate(rows):
    y = 1.96 + idx * 0.56
    fill = WHITE if idx % 2 == 0 else "EDF2F6"
    box(slide, 0.65, y, 4.35, 0.50, row[0], fill=fill, line=LINE, radius=False,
        size=11, align=PP_ALIGN.LEFT)
    for col, enabled, color in [(5.02, row[1], GREEN), (6.74, row[2], ORANGE), (8.10, row[3], PURPLE)]:
        box(slide, col, y, 1.68 if col == 5.02 else (1.32 if col == 6.74 else 1.46), 0.50,
            "●" if enabled else "—", fill=fill, line=LINE, radius=False,
            size=15, color=color if enabled else "A9B5C0", bold=enabled)
    box(slide, 9.60, y, 3.05, 0.50, row[4], fill=fill, line=LINE, radius=False,
        size=10.2, align=PP_ALIGN.LEFT)
status_legend(slide)


# 5 — expected orchestration
slide = slide_frame(
    "预期编排：七个 Agent 节点全部保留，评测与发布 Gate 不是 Agent",
    "IMPLEMENTED + PLANNED",
    "代码依据：infrastructure/domain-knowledge/src/graph.ts:13-27,106-188；domain-knowledge/docs/report/01-Agent输入输出总览.md",
    dark=True,
    kicker="ORCHESTRATION",
)
node(slide, 0.55, 2.88, 1.40, 0.80, "Orchestrator", "拆任务 / 组装上下文", fill=NAVY_2,
     accent=CYAN, title_color=WHITE, body_color="B9C9D8")
chevron(slide, 2.02, 3.06, 0.26, 0.40, CYAN)
node(slide, 2.35, 1.65, 1.40, 0.80, "DocWorker × N", "分块证据片段", fill=NAVY_2,
     accent=BLUE, title_color=WHITE, body_color="B9C9D8")
node(slide, 2.35, 4.18, 1.40, 0.80, "TestGen", "oracle / 测试意图", fill=NAVY_2,
     accent=PURPLE, title_color=WHITE, body_color="B9C9D8")
node(slide, 4.20, 1.65, 1.35, 0.80, "DocGen", "候选知识", fill=NAVY_2,
     accent=BLUE, title_color=WHITE, body_color="B9C9D8")
node(slide, 4.20, 4.18, 1.35, 0.80, "Oracle 校验", "非 Agent", fill=NAVY,
     accent=ORANGE, title_color=WHITE, body_color="B9C9D8")
node(slide, 6.05, 1.65, 1.25, 0.80, "Code", "fresh 实现", fill=NAVY_2,
     accent=GREEN, title_color=WHITE, body_color="B9C9D8")
node(slide, 7.78, 1.65, 1.25, 0.80, "Check", "只读检查", fill=NAVY_2,
     accent=ORANGE, title_color=WHITE, body_color="B9C9D8")
node(slide, 7.04, 3.52, 1.65, 0.80, "EvalRunner", "独立确定性评测\n非 Agent", fill=NAVY,
     accent=RED, title_color=WHITE, body_color="B9C9D8")
node(slide, 9.48, 3.52, 1.35, 0.80, "Review", "归因 / Correction", fill=NAVY_2,
     accent=PURPLE, title_color=WHITE, body_color="B9C9D8")
node(slide, 11.28, 3.52, 1.42, 0.80, "Publication Gate", "确定性策略\n非 Agent", fill=NAVY,
     accent=GREEN, title_color=WHITE, body_color="B9C9D8")
for x, y, color in [(3.82, 1.84, BLUE), (5.66, 1.84, GREEN), (7.40, 1.84, ORANGE),
                    (3.82, 4.37, PURPLE), (8.85, 3.72, PURPLE), (10.92, 3.72, GREEN)]:
    chevron(slide, x, y, 0.25, 0.37, color)
line(slide, 8.38, 2.52, 8.02, 3.46, ORANGE, 2)
line(slide, 4.88, 5.05, 7.47, 4.39, PURPLE, 2)
line(slide, 12.02, 4.40, 12.02, 5.63, ORANGE, 2)
line(slide, 12.02, 5.63, 1.25, 5.63, ORANGE, 2)
line(slide, 1.25, 5.63, 1.25, 3.77, ORANGE, 2)
textbox(slide, 4.30, 5.73, 4.1, 0.30, "ITERATE / ROLLBACK：回到新一轮 Orchestrator", 10.5, ORANGE, True,
        align=PP_ALIGN.CENTER)
status_legend(slide, dark=True)


# 6 — I/O overview
slide = slide_frame(
    "七个 Agent 输入输出：职责没有简化，交接只通过受约束产物",
    "IMPLEMENTED + PLANNED",
    "对照：contracts/src/index.ts:200-216；infrastructure/domain-knowledge/src/agent-definitions.ts；domain-knowledge 报告 §Agent 总表",
    kicker="AGENT CONTRACTS",
)
cols = [("Agent", 0.58, 1.42), ("主要输入", 2.00, 4.05), ("主要输出", 6.05, 4.25), ("边界", 10.30, 2.43)]
for label, x, w in cols:
    box(slide, x, 1.43, w, 0.44, label, fill=NAVY, line=NAVY, color=WHITE, size=10.5, bold=True)
agent_rows = [
    ("Orchestrator", "RunPolicy、snapshot、历史 Correction", "任务 DAG、worker 划分、上下文计划", "不生成最终知识"),
    ("DocWorker", "语义块、依赖摘要、公开接口", "证据片段、provenance、未决问题", "不跨块猜测"),
    ("DocGen", "全部 worker 片段、写作规范、历史版本", "候选知识 + 事实引用", "不发布"),
    ("TestGen", "知识契约、公开接口、风险项", "oracle 意图、候选命令/用例", "不读取候选实现"),
    ("Code", "候选知识、接口约束、allowed paths", "fresh generated files", "不读取参考实现"),
    ("Check", "实现 diff、规则、静态证据", "blocking + findings + 位置", "只读、不定发布"),
    ("Review", "Eval、Check、候选知识、历史", "归因、Correction、unresolved risks", "只读、不选最终状态"),
]
for idx, row in enumerate(agent_rows):
    y = 1.92 + idx * 0.64
    fill = WHITE if idx % 2 == 0 else "EDF2F6"
    widths = [1.42, 4.05, 4.25, 2.43]
    xs = [0.58, 2.00, 6.05, 10.30]
    for col, (value, x, w) in enumerate(zip(row, xs, widths)):
        box(slide, x, y, w, 0.58, value, fill=fill, line=LINE, radius=False,
            size=10.5 if col else 11, color=NAVY if col == 0 else INK,
            bold=col == 0, align=PP_ALIGN.LEFT if col else PP_ALIGN.CENTER)
status_legend(slide)


# 7 — doc current
slide = slide_frame(
    "文档生成阶段：并行结构已跑通，但当前内容仍是 deterministic fixture",
    "DEMO",
    "代码依据：graph.ts:137-153；automated-project-workflow.ts:143-182；tests/integration/langgraph-infrastructure.test.ts",
    kicker="DOCUMENTATION · CURRENT",
)
pill(slide, 0.62, 1.45, 1.15, "当前可运行", ORANGE, size=9)
node(slide, 0.62, 2.23, 1.45, 0.84, "Orchestrator", "workerCount", fill=WHITE, accent=CYAN)
chevron(slide, 2.16, 2.43, color=CYAN)
for i in range(3):
    node(slide, 2.58, 1.54 + i * 1.12, 1.62, 0.74, f"DocWorker {i + 1}",
         "workerId + 固定 fragment", fill=BLUE_LIGHT, accent=BLUE)
line(slide, 4.30, 1.91, 5.20, 2.62, BLUE, 1.5)
line(slide, 4.30, 3.03, 5.20, 2.62, BLUE, 1.5)
line(slide, 4.30, 4.15, 5.20, 2.62, BLUE, 1.5)
node(slide, 5.24, 2.18, 1.72, 0.88, "DocGen", "读取预置 knowledge-v1/v2", fill=BLUE_LIGHT, accent=BLUE)
chevron(slide, 7.05, 2.41, color=CYAN)
node(slide, 7.45, 2.18, 1.72, 0.88, "Candidate", "CAS + provenance", fill=GREEN_LIGHT, accent=GREEN)

rich_box(slide, 9.55, 1.46, 3.15, 2.02, "这个 demo 能证明", [
    "Send fan-out / join 路径可执行",
    "多个 worker ArtifactRef 能汇入 DocGen context",
    "节点投影、checkpoint、候选写入可观察",
], accent=GREEN)
rich_box(slide, 9.55, 3.73, 3.15, 2.02, "它不能证明", [
    "没有函数 / 类 / SCC 级语义切块",
    "DocGen 没有按片段内容聚合写作",
    "没有增量上下文、冲突消解和覆盖率验证",
], accent=RED)
box(slide, 0.62, 5.22, 8.55, 0.54,
    "结论：当前是“拓扑形状正确”的 fixture，不是“拓扑与语义切块已经实现”。",
    fill=ORANGE_LIGHT, line=ORANGE, color=INK, size=13, bold=True)
status_legend(slide)


# 8 — semantic chunking target
slide = slide_frame(
    "语义 / 拓扑切块：以可解释代码单元为边界，不按 token 生硬截断",
    "PLANNED",
    "研究依据：DocAgent, arXiv:2504.08725 (2025-04-11)；domain-knowledge/docs/report/01-Agent输入输出总览.md §文档阶段",
    dark=True,
    kicker="DOCUMENTATION · TARGET",
)
textbox(slide, 0.65, 1.42, 4.0, 0.30, "1 · 建图与稳定边界", 12, CYAN, True)
graph_nodes = [
    (0.72, 2.05, "public API", BLUE), (2.45, 1.68, "service", PURPLE),
    (2.45, 2.72, "repository", GREEN), (4.18, 2.20, "adapter", ORANGE),
]
for x, y, label, color in graph_nodes:
    node(slide, x, y, 1.20, 0.60, label, "", fill=NAVY_2, accent=color,
         title_color=WHITE, body_color="B9C9D8")
line(slide, 1.92, 2.34, 2.42, 2.02, CYAN, 1.5)
line(slide, 1.92, 2.34, 2.42, 3.02, CYAN, 1.5)
line(slide, 3.67, 2.02, 4.15, 2.46, CYAN, 1.5)
line(slide, 3.67, 3.02, 4.15, 2.46, CYAN, 1.5)
textbox(slide, 0.72, 3.62, 4.72, 0.66,
        "AST / symbol / import-call graph → SCC 收缩 → 拓扑层\n函数、类、模块、测试簇是候选块；超预算时只在内部语义边界再切。",
        11.5, "C5D3E0", valign=MSO_ANCHOR.TOP)

textbox(slide, 5.80, 1.42, 3.6, 0.30, "2 · DocWorker 产出证据包", 12, CYAN, True)
rich_box(slide, 5.78, 1.83, 3.02, 2.57, "ChunkEvidence", [
    "chunkId / symbols / source ranges",
    "imports / callers / callees / tests",
    "事实、约束、示例、provenance",
    "依赖摘要与 unresolved issues",
], fill=NAVY_2, accent=BLUE, title_color=WHITE, body_color="C5D3E0")

textbox(slide, 9.18, 1.42, 3.4, 0.30, "3 · DocGen 增量聚合", 12, CYAN, True)
for i, (label, sub, color) in enumerate([
    ("按拓扑序消费", "先依赖，后调用方", BLUE),
    ("增量上下文", "摘要可复用，正文不重复", PURPLE),
    ("验证—重写闭环", "覆盖 / 引用 / 冲突检查", GREEN),
]):
    node(slide, 9.18, 1.84 + i * 0.91, 3.42, 0.67, label, sub, fill=NAVY_2,
         accent=color, title_color=WHITE, body_color="C5D3E0")
box(slide, 0.67, 4.78, 11.93, 1.00, fill="173A5E", line="31516E", text="")
textbox(slide, 0.93, 4.96, 2.18, 0.26, "避免语义截断的判据", 11, CYAN, True)
textbox(slide, 3.00, 4.91, 9.22, 0.58,
        "块内必须可独立解释；跨块依赖显式引用；声明与测试共同归组；循环依赖作为 SCC 整体；\n超 token 时保留接口 + 摘要 + provenance，正文递延，禁止从任意字符位置截断。",
        12.2, WHITE, True, valign=MSO_ANCHOR.TOP)
textbox(slide, 0.68, 6.16, 11.9, 0.28,
        "DocAgent 论文验证了 topological code processing + incremental context building；本项目采用其原则，具体 Schema 与 Gate 仍需在 wpKnowledge 内实现。",
        10.5, "B9C9D8", italic=True)
status_legend(slide, dark=True)


# 9 — code and eval
slide = slide_frame(
    "代码生成与独立评测：路径隔离已存在，live CodeAgent 尚未接线",
    "IMPLEMENTED + DEMO",
    "代码依据：automated-project-workflow.ts:229-309；contracts/src/index.ts:127-197；composition.ts:100-109",
    kicker="CODE + EVALUATION",
)
textbox(slide, 0.63, 1.43, 3.3, 0.28, "候选路径 · 生成后再检查", 11, BLUE, True)
node(slide, 0.63, 2.05, 1.40, 0.78, "DocGen", "候选知识", fill=BLUE_LIGHT, accent=BLUE)
chevron(slide, 2.12, 2.24, color=BLUE)
node(slide, 2.55, 2.05, 1.40, 0.78, "Code", "静态 code-v1/v2", fill=GREEN_LIGHT, accent=GREEN)
chevron(slide, 4.04, 2.24, color=GREEN)
node(slide, 4.46, 2.05, 1.40, 0.78, "Check", "blocking / findings", fill=ORANGE_LIGHT, accent=ORANGE)

textbox(slide, 0.63, 3.49, 3.3, 0.28, "Oracle 路径 · 与候选实现隔离", 11, PURPLE, True)
node(slide, 0.63, 4.10, 1.40, 0.78, "TestGen", "固定命令 fixture", fill=PURPLE_LIGHT, accent=PURPLE)
chevron(slide, 2.12, 4.29, color=PURPLE)
node(slide, 2.55, 4.10, 1.82, 0.78, "Reference Oracle", "先验证测试基线", fill=PURPLE_LIGHT, accent=PURPLE)

line(slide, 5.88, 2.45, 6.62, 3.23, ORANGE, 2)
line(slide, 4.40, 4.49, 6.62, 3.23, PURPLE, 2)
node(slide, 6.68, 2.73, 2.10, 1.02, "TrustedProjectEvaluator", "归档 workspace 内执行命令\n采集 stdout / stderr / tests", fill=NAVY, accent=RED,
     title_color=WHITE, body_color="D2DDE7")
chevron(slide, 8.90, 3.03, color=RED)
node(slide, 9.38, 2.73, 1.63, 1.02, "Evaluation", "不可变 evidenceRef", fill=RED_LIGHT, accent=RED)
chevron(slide, 11.11, 3.03, color=RED)
node(slide, 11.56, 2.73, 1.15, 1.02, "Review", "归因", fill=PURPLE_LIGHT, accent=PURPLE)

rich_box(slide, 6.68, 4.30, 2.90, 1.47, "已实现", [
    "allowed path 校验",
    "reference oracle 先行",
    "真实 Node 命令 + CAS 证据",
], accent=GREEN)
rich_box(slide, 9.81, 4.30, 2.90, 1.47, "当前边界", [
    "Code / TestGen 仍读 fixture",
    "不是 CodeAgent provider",
    "受信执行器不是敌对沙箱",
], accent=ORANGE)
status_legend(slide)


# 10 — review/gate
slide = slide_frame(
    "Review / Correction / 迭代：Review 给证据，Gate 才给状态",
    "IMPLEMENTED",
    "代码依据：automated-project-workflow.ts:312-374；application/src/index.ts:128-194；domain/src/index.ts:159-196；commit b6973a8",
    dark=True,
    kicker="REVIEW LOOP",
)
node(slide, 0.63, 2.26, 1.58, 0.92, "Evaluation", "tests / stability\ninfrastructureFailure", fill=NAVY_2,
     accent=RED, title_color=WHITE, body_color="C5D3E0")
chevron(slide, 2.31, 2.51, color=RED)
node(slide, 2.72, 2.26, 1.58, 0.92, "Review", "blocking / recommendation\nCorrection", fill=NAVY_2,
     accent=PURPLE, title_color=WHITE, body_color="C5D3E0")
chevron(slide, 4.41, 2.51, color=PURPLE)
node(slide, 4.82, 2.07, 2.04, 1.30, "wp Domain Gate", "绑定 body / code / oracle / check / review\n策略产生唯一 GateDecision", fill="224867",
     accent=GREEN, title_color=WHITE, body_color="D4E0EA")

for i, (label, sub, color) in enumerate([
    ("PASS", "Publication → VERIFIED", GREEN),
    ("ITERATE", "FlywheelRun +1 → 新候选", ORANGE),
    ("STOPPED", "LOW_CONFIDENCE", RED),
]):
    y = 1.47 + i * 1.16
    chevron(slide, 7.05, y + 0.22, color=color)
    node(slide, 7.43, y, 2.02, 0.78, label, sub, fill=NAVY_2,
         accent=color, title_color=WHITE, body_color="C5D3E0")

rich_box(slide, 9.92, 1.47, 2.75, 2.88, "修复后的约束", [
    "Review 产物必须先存在，才记录正常 GateDecision",
    "Check / Review blocking 进入 reasonCodes",
    "Graph 只执行 GateDecision，不再另算 maxIterations",
    "基础设施失败可跳过 Review，但仍由 Domain Gate STOPPED",
], fill=NAVY_2, accent=CYAN, title_color=WHITE, body_color="C5D3E0")
box(slide, 0.65, 4.75, 8.78, 0.84, fill="173A5E", line="31516E", text="")
textbox(slide, 0.91, 4.94, 8.25, 0.40,
        "Correction 只描述“知识路径 → 失败证据 → 可判定修订标准”；\n它不能直接改状态，也不能绕过 fresh Code / Eval / Gate。",
        13, WHITE, True, align=PP_ALIGN.CENTER)
status_legend(slide, dark=True)


# 11 — dual-lane end-to-end
slide = slide_frame(
    "端到端流程：LangGraph 推进执行，wpKnowledge 记录可发布事实",
    "IMPLEMENTED + DEMO",
    "代码依据：graph.ts；runtime.ts；automated-project-workflow.ts；application/src/index.ts；SQLiteFlywheelRepository",
    kicker="END TO END",
)
textbox(slide, 0.60, 1.42, 1.70, 0.35, "LANGGRAPH", 12, BLUE, True)
textbox(slide, 0.60, 4.15, 1.70, 0.35, "WPKNOWLEDGE", 12, GREEN, True)
line(slide, 0.58, 3.77, 12.72, 3.77, LINE, 1.2)
top = [
    ("Orchestrate", CYAN), ("DocWorker\n+ DocGen", BLUE), ("Code + Check", GREEN),
    ("Eval", RED), ("Review", PURPLE), ("Router", ORANGE), ("Publication", GREEN),
]
xs = [0.62, 2.25, 4.02, 5.79, 7.33, 8.86, 10.55]
ws = [1.28, 1.42, 1.42, 1.20, 1.20, 1.20, 1.52]
for i, ((label, color), x, w) in enumerate(zip(top, xs, ws)):
    box(slide, x, 2.10, w, 0.72, label, fill=WHITE, line=color, size=11.5, bold=True)
    if i < len(top) - 1:
        chevron(slide, x + w + 0.08, 2.28, 0.22, 0.34, color)
bottom = [
    ("CREATED\n→ GENERATING", BLUE), ("CANDIDATE\n+ CAS refs", BLUE),
    ("GenerationKey\nCOMMITTED", PURPLE), ("EvaluationReport\n+ evidence", RED),
    ("GateDecision", ORANGE), ("PUBLISHING", GREEN), ("VERIFIED", GREEN),
]
for i, ((label, color), x, w) in enumerate(zip(bottom, xs, ws)):
    box(slide, x, 4.72, w, 0.72, label, fill=WHITE, line=color, size=10.5, bold=True)
    if i < len(bottom) - 1:
        chevron(slide, x + w + 0.08, 4.90, 0.22, 0.34, color)
for x in [1.26, 2.96, 4.73, 6.39, 7.93, 9.46, 11.31]:
    line(slide, x, 2.88, x, 4.66, "A9BAC8", 1.0)
box(slide, 9.87, 5.82, 2.85, 0.52,
    "唯一可产生 VERIFIED 的入口", fill=GREEN, line=GREEN, color=WHITE, size=12.2, bold=True)
line(slide, 11.31, 5.46, 11.31, 5.81, GREEN, 2.2)
textbox(slide, 0.65, 5.84, 8.55, 0.47,
        "ITERATE：GateDecision → FlywheelRun.ITERATING → Graph iteration + 1 → 新一轮；\nGraph checkpoint 不能替代业务事件，业务事件也不能恢复节点局部执行。",
        11.2, MUTED, True, valign=MSO_ANCHOR.TOP)
status_legend(slide)


# 12 — demo limits
slide = slide_frame(
    "ohMyWorkPanel deterministic fixture：可证明路径，不可证明智能能力",
    "DEMO",
    "代码依据：acceptance/ohmyworkpanel/*；tests/acceptance/automated-langgraph-flow.test.ts；automated-project-workflow.ts:137-182",
    kicker="DEMO EVIDENCE",
)
box(slide, 0.64, 1.48, 5.78, 0.46, "可以据此下结论", fill=GREEN, line=GREEN, color=WHITE, size=13, bold=True)
box(slide, 6.88, 1.48, 5.78, 0.46, "不能据此下结论", fill=RED, line=RED, color=WHITE, size=13, bold=True)
proofs = [
    "七个 Agent 节点均经过 LangGraph",
    "DocWorker fan-out 后，多个引用到达 DocGen context",
    "第一轮失败 → ITERATE → 第二轮 PASS",
    "Oracle / Check / Review 引用进入最终 EvaluationReport",
    "只有 wpKnowledge Publication 将版本置为 VERIFIED",
]
limits = [
    "Agent 会基于代码自主生成高质量内容",
    "DocWorker 已按 AST / 拓扑 / 语义完成分块",
    "promptAddon 已影响当前 fixture 的实际输出",
    "worker 真并行吞吐、乱序聚合与部分失败恢复已可靠",
    "TrustedProjectEvaluator 等同于不可信代码沙箱",
]
for i, item in enumerate(proofs):
    box(slide, 0.64, 2.12 + i * 0.76, 5.78, 0.60, f"✓  {item}", fill=GREEN_LIGHT,
        line="B8DEC9", size=11.5, color=INK, bold=i in (2, 4), align=PP_ALIGN.LEFT)
for i, item in enumerate(limits):
    box(slide, 6.88, 2.12 + i * 0.76, 5.78, 0.60, f"×  {item}", fill=RED_LIGHT,
        line="EAC1BE", size=11.5, color=INK, bold=i in (1, 4), align=PP_ALIGN.LEFT)
status_legend(slide)


# 13 — local issues
slide = slide_frame(
    "本地运行记录：成功项与环境问题分开描述",
    "IMPLEMENTED + DEMO",
    "核对依据：package.json / .github/workflows/ci.yml / package-lock.json；domain-knowledge/src/graph/build-graph.ts；本次本机实测",
    dark=True,
    kicker="RUNTIME NOTES",
)
headers = [("事项", 0.62, 2.45), ("实际观察", 3.07, 5.22), ("处理 / 当前口径", 8.29, 4.40)]
for label, x, w in headers:
    box(slide, x, 1.46, w, 0.43, label, fill="224867", line="31516E", color=WHITE, size=10.5, bold=True)
runtime_rows = [
    ("Node 版本", "本机 Node 22.22；仓库 engines 与 CI 为 Node 24", "本地可跑；正式结果以 Node 24 CI 为准"),
    ("node:sqlite", "Node 22 测试输出 experimental warning", "属于运行时提示，不是测试失败"),
    ("lockfile registry", "首次写入腾讯镜像，GitHub CI npm ci 出现 ENOTFOUND", "776dac5 已改为 registry.npmjs.org"),
    ("CodeAgent timeout", "domain-knowledge doc-gen 实测约 2–3 分钟，原节点 120s", "上调到 600s；wpKnowledge graph 同为 600s"),
    ("真实 Agent 接入", "wp 自动路径实例化 OhMyWorkPanel fixture executor", "AgentProvider / CodeAgent provider 尚未接入"),
    ("本次验证", "typecheck、validate:specs、npm test", "全部通过；50 / 50 tests"),
]
for idx, row in enumerate(runtime_rows):
    y = 1.95 + idx * 0.70
    fill = NAVY_2 if idx % 2 == 0 else "1B4265"
    for value, x, w in zip(row, [0.62, 3.07, 8.29], [2.45, 5.22, 4.40]):
        box(slide, x, y, w, 0.63, value, fill=fill, line="31516E", radius=False,
            size=10.7, color=WHITE if x == 0.62 else "D6E1EA",
            bold=x == 0.62, align=PP_ALIGN.LEFT)
status_legend(slide, dark=True)


# 14 — next steps
slide = slide_frame(
    "下一步：先把语义证据链做实，再扩大 Agent 与运行时能力",
    "PLANNED",
    "评审结论汇总：PR #20 @ b6973a8；domain-knowledge 对照；DocAgent arXiv:2504.08725",
    kicker="NEXT STEPS",
)
priorities = [
    ("P0", "语义 / 拓扑分块", "构建 symbol + import/call/test 图；SCC 与拓扑层；定义 ChunkEvidence Schema；DocGen 真正消费所有 worker 产物。", RED),
    ("P0", "接入真实 Agent provider", "让 promptAddon 进入真实请求；固定 base prompt / tools / schema；输出做运行时校验，禁止 Agent 直接决定发布。", RED),
    ("P1", "并行、恢复、取消", "补真实 overlap、乱序聚合、worker 部分失败；清理 RUNNING GenerationKey；对子进程做超时与强制终止。", ORANGE),
    ("P1", "评测隔离", "从 TrustedProjectEvaluator 演进到不可信代码沙箱；收紧网络、文件、CPU / 内存 / 输出预算。", ORANGE),
    ("P2", "增量与质量闭环", "按变更影响子图复用上下文；增加覆盖、引用、冲突、一致性 Gate；用真实仓库样本建立回归集。", PURPLE),
]
for idx, (priority, title, detail, color) in enumerate(priorities):
    y = 1.46 + idx * 0.95
    pill(slide, 0.63, y + 0.14, 0.58, priority, color, size=9)
    box(slide, 1.35, y, 11.35, 0.78, fill=WHITE, line=LINE, text="")
    box(slide, 1.35, y, 0.07, 0.78, fill=color, line=color, radius=False, text="")
    textbox(slide, 1.63, y + 0.10, 2.60, 0.26, title, 13.5, NAVY, True)
    textbox(slide, 4.12, y + 0.09, 8.28, 0.50, detail, 11.3, MUTED,
            valign=MSO_ANCHOR.TOP)
box(slide, 0.63, 6.28, 12.07, 0.43,
    "验收原则：每项能力必须用真实输入、不可变证据和失败用例证明；计划项不得借 fixture 提前标记为 implemented。",
    fill=NAVY, line=NAVY, color=WHITE, size=12, bold=True)
status_legend(slide)


OUTPUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUTPUT)
print(f"generated {OUTPUT} slides={len(prs.slides)}")
